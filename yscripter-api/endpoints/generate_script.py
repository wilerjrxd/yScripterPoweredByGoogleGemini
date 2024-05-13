from flask import Blueprint, request, jsonify
from PIL import Image
from IPython.display import Markdown
import re
import tempfile
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PyPDF2 import PdfReader



import google.generativeai as genai
from google.cloud import documentai_v1 as documentai
import os
import json

generate_script_bp = Blueprint('post_image', __name__)

UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'mp4'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def to_markdown(text):
  text = text.replace('•', '  *')
  return Markdown(text.indent(text, '> ', predicate=lambda _: True))

@generate_script_bp.route('/generate-script', methods=['POST'])
def generateScript():

    theme = request.form.get('theme')
    style = request.form.get('style')
    reutilize = request.form.get('reutilize')
    details = request.form.get('details')
    duration = request.form.get('duration')

    files = request.files.getlist('file[]')
    media = 'None'

    API_KEY = <GOOGLE_API_KEY>

    genai.configure(api_key=API_KEY)

    generation_config = {
        "temperature": float(style),
        "max_output_tokens": 500
    }

    model = genai.GenerativeModel(
        model_name='gemini-1.5-pro-latest',
        generation_config=generation_config,
    )

    safety_settings = [
        {
            "category": "HARM_CATEGORY_HARASSMENT",
            "threshold": "BLOCK_ONLY_HIGH"
        },
        {
            "category": "HARM_CATEGORY_HATE_SPEECH",
            "threshold": "BLOCK_ONLY_HIGH"
        },
        {
            "category": "HARM_CATEGORY_SEXUALLY_EXPLICIT",
            "threshold": "BLOCK_ONLY_HIGH"
        },
        {
            "category": "HARM_CATEGORY_DANGEROUS_CONTENT",
            "threshold": "BLOCK_ONLY_HIGH"
        },
    ]

    file_path = ""

    prompt = "aja como um roteirista profissional. Escreva o texto para um vídeo com duração " + str(duration) + " minutos considerando como parâmetros: tema: " + theme
    if (reutilize != 'undefined'):
        prompt += "; com base na estrutura do texto: " + reutilize
    if (details != 'undefined'):
        prompt += "; detalhes adicionais: " + details

    response = ""

    if len(files) > 0:
        for file in files:
            file_dir = tempfile.gettempdir()
            file_path = os.path.join(file_dir, file.filename)
            file.save(file_path)
            media = genai.upload_file(path=file_path, display_name="Arquivo")
            if (is_document(file_path) == False):
                prompt += "; com base no conteúdo da mídia em anexo"
                response = model.generate_content([prompt, media], stream=True)
            else:
                media = extract_text_from_document(file_path)
                prompt += "; com base no conteúdo extraido de um documento: " + media
                response = model.generate_content([prompt], stream=True)

    response.resolve()

    resultText = ''

    if (len(response.parts) > 0):
        for part in response.parts:
            resultText += part.text
    else:
        resultText = response.text

    resultText = resultText.split('\n\n')

    os.remove(file_path)

    return jsonify(resultText), 200

def is_document(filepath):
    document_extensions = ['.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.txt']
    extension = get_last_four_digits(filepath)
    return extension in document_extensions

def extract_text_from_document(document_path):
    if get_last_four_digits(document_path) in ['doc', 'docx']:
        print('doc')
        return extract_text_from_docx(document_path)
    elif get_last_four_digits(document_path) in 'pdf':
        print('pdf')

        return extract_text_from_pdf(document_path)
    elif get_last_four_digits(document_path) in ['xls', 'xlsx']:
        print('xls')
        return extract_text_from_xlsx(document_path)
    elif get_last_four_digits(document_path) in ['ppt', 'pptx']:
        print('ppt')
        return extract_text_from_pptx(document_path)
    elif get_last_four_digits(document_path) in 'txt':
        print('txt')
        return extract_text_from_txt(document_path)
    else:
        return 'None'
    

def extract_text_from_txt(txt_path):
    text = ''
    with open(txt_path, 'r', encoding='utf-8') as file:
        for line in file:
            text += line
    return text

def extract_text_from_pdf(pdf_path):
    with open(pdf_path, 'rb') as pdf_file:
        reader = PdfReader(pdf_file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
        return text

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = '\n'.join(paragraph.text for paragraph in doc.paragraphs)
    return text

def extract_text_from_xlsx(xlsx_path):
    wb = load_workbook(xlsx_path)
    text = ''
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join(str(cell) for cell in row if cell) + '\n'
    return text


def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + '\n'
    return text

def get_last_four_digits(string):
    return string[-4:]
